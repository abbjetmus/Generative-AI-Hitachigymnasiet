"""Script to generate the Generativ AI introductory PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy
BG_GRADIENT = RGBColor(0x1A, 0x1A, 0x2E)    # Dark purple-navy
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)    # Bright blue
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)  # Purple
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)    # Cyan
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # Green
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)  # Orange
ACCENT_PINK = RGBColor(0xEC, 0x48, 0x99)    # Pink
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MEDIUM_GRAY = RGBColor(0x94, 0xA3, 0xB8)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_icon_text(slide, left, top, icon_text, label, icon_color, desc="", desc_color=LIGHT_GRAY):
    """Add an icon (emoji/symbol in a circle) with label text."""
    # Icon circle
    circle = add_circle(slide, left, top, Inches(0.8), icon_color)
    circle.text_frame.paragraphs[0].text = icon_text
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.bold = True

    # Label
    add_text_box(slide, left - Inches(0.3), top + Inches(0.9), Inches(1.4), Inches(0.5),
                 label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if desc:
        add_text_box(slide, left - Inches(0.3), top + Inches(1.3), Inches(1.4), Inches(0.8),
                     desc, font_size=11, color=desc_color, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Decorative shapes
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x7C, 0x3A, 0xED))
add_circle(slide, Inches(10), Inches(4), Inches(5), RGBColor(0x00, 0x96, 0xFF))

# Dark overlay
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
from lxml import etree
a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
spPr = overlay._element.find(f'.//{{{a_ns}}}solidFill/{{{a_ns}}}srgbClr')
if spPr is not None:
    alpha = etree.SubElement(spPr, f'{{{a_ns}}}alpha')
    alpha.set('val', '75000')
overlay.line.fill.background()

# Title text
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "GENERATIV AI", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "En introduktion till framtidens teknologi", font_size=28, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.8), Inches(3.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Course info
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "Hitachigymnasiet", font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
             "Kursmaterial: github.com/microsoft/generative-ai-with-javascript",
             font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Vad är AI?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Vad är Artificiell Intelligens?", font_size=40, color=WHITE, bold=True)

# Left side - description
desc_box = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), RGBColor(0x1E, 0x29, 0x3B), ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(4.6),
             "AI (Artificiell Intelligens) är en samling verktyg och tekniker som gör det möjligt för datorer att utföra uppgifter som normalt kräver mänsklig intelligens.\n\n"
             "AI finns redan överallt i våra liv:\n\n"
             "- Varje gång du googlar något\n"
             "- När du betalar med kort (bedrägeridetektering)\n"
             "- Netflix och Spotify rekommendationer\n"
             "- Autokorrigering på din telefon\n"
             "- Ansiktsigenkänning\n\n"
             "Du använder AI dussintals gånger om dagen - ofta utan att tänka på det!",
             font_size=16, color=LIGHT_GRAY)

# Right side - AI types
box1 = add_shape(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(1.4), RGBColor(0x1E, 0x29, 0x3B), ACCENT_PURPLE)
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.4),
             "Supervised Learning", font_size=20, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.7),
             "Lär sig mönster från märkt data.\nT.ex. spamfilter, bildklassificering.",
             font_size=14, color=LIGHT_GRAY)

box2 = add_shape(slide, Inches(7), Inches(3.3), Inches(5.5), Inches(1.4), RGBColor(0x1E, 0x29, 0x3B), ACCENT_CYAN)
add_text_box(slide, Inches(7.3), Inches(3.4), Inches(5), Inches(0.4),
             "Generativ AI", font_size=20, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(7.3), Inches(3.9), Inches(5), Inches(0.7),
             "Skapar nytt innehåll: text, bilder, ljud, kod.\nT.ex. ChatGPT, DALL-E, GitHub Copilot.",
             font_size=14, color=LIGHT_GRAY)

box3 = add_shape(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(1.4), RGBColor(0x1E, 0x29, 0x3B), ACCENT_GREEN)
add_text_box(slide, Inches(7.3), Inches(5.1), Inches(5), Inches(0.4),
             "Reinforcement Learning", font_size=20, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(7.3), Inches(5.6), Inches(5), Inches(0.7),
             "Lär sig genom trial and error.\nT.ex. spel-AI, robotik.",
             font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 3: Vad är Generativ AI?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Vad är Generativ AI?", font_size=40, color=WHITE, bold=True)

# Main definition box
def_box = add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.5), RGBColor(0x1E, 0x29, 0x3B), ACCENT_CYAN)
add_text_box(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(1.3),
             "Generativ AI är AI-system som kan skapa nytt, högkvalitativt innehåll - text, bilder, ljud och kod.\n"
             "Det mest kända exemplet är ChatGPT, som lanserades i november 2022 och förändrade allt.",
             font_size=20, color=WHITE)

# Impact stats
stat_y = Inches(3.4)
stat_h = Inches(3.2)

s1 = add_shape(slide, Inches(0.8), stat_y, Inches(3.6), stat_h, RGBColor(0x1E, 0x29, 0x3B), ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(3.6), Inches(3), Inches(0.6),
             "$2.6-4.4", font_size=36, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.1), Inches(4.3), Inches(3), Inches(0.4),
             "BILJONER USD", font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.1), Inches(4.8), Inches(3), Inches(1.5),
             "Beräknat årligt tillskott till den globala ekonomin tack vare generativ AI",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

s2 = add_shape(slide, Inches(4.9), stat_y, Inches(3.6), stat_h, RGBColor(0x1E, 0x29, 0x3B), ACCENT_PURPLE)
add_text_box(slide, Inches(5.2), Inches(3.6), Inches(3), Inches(0.6),
             "7%", font_size=36, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.2), Inches(4.3), Inches(3), Inches(0.4),
             "GLOBAL BNP-ÖKNING", font_size=14, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.2), Inches(4.8), Inches(3), Inches(1.5),
             "Goldman Sachs uppskattar att generativ AI kan höja global BNP med 7% under nästa decennium",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

s3 = add_shape(slide, Inches(9), stat_y, Inches(3.6), stat_h, RGBColor(0x1E, 0x29, 0x3B), ACCENT_GREEN)
add_text_box(slide, Inches(9.3), Inches(3.6), Inches(3), Inches(0.6),
             "80%", font_size=36, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.3), Inches(4.3), Inches(3), Inches(0.4),
             "AV ALLA ARBETARE", font_size=14, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.3), Inches(4.8), Inches(3), Inches(1.5),
             "Över 80% av arbetare kan få minst 10% av sina arbetsuppgifter påverkade av generativ AI",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: Hur fungerar LLMs?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Hur fungerar stora språkmodeller (LLMs)?", font_size=38, color=WHITE, bold=True)

# Step by step explanation
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "LLMs lär sig att förutsäga nästa ord - och det är hela hemligheten!", font_size=20, color=ACCENT_CYAN)

# Example box
ex_box = add_shape(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B), ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(11), Inches(0.4),
             "Exempel: Hur modellen lär sig", font_size=20, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.9), Inches(11), Inches(1.4),
             'Mening: "Min favoritmat är en bagel med cream cheese"\n\n'
             '"Min favoritmat är en"  -->  bagel\n'
             '"Min favoritmat är en bagel"  -->  med\n'
             '"Min favoritmat är en bagel med"  -->  cream\n'
             '"Min favoritmat är en bagel med cream"  -->  cheese',
             font_size=16, color=LIGHT_GRAY, font_name="Consolas")

# Key facts
facts_y = Inches(4.7)

f1 = add_shape(slide, Inches(0.8), facts_y, Inches(3.6), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B), ACCENT_PURPLE)
add_text_box(slide, Inches(1.1), Inches(4.85), Inches(3), Inches(0.4),
             "Massiv träningsdata", font_size=18, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(1.1), Inches(5.35), Inches(3), Inches(1.2),
             "Hundratals miljarder till biljoner ord från internet, böcker och andra källor",
             font_size=14, color=LIGHT_GRAY)

f2 = add_shape(slide, Inches(4.9), facts_y, Inches(3.6), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B), ACCENT_CYAN)
add_text_box(slide, Inches(5.2), Inches(4.85), Inches(3), Inches(0.4),
             "Enorma modeller", font_size=18, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(5.2), Inches(5.35), Inches(3), Inches(1.2),
             "Miljarder parametrar som fångar språkets mönster och struktur",
             font_size=14, color=LIGHT_GRAY)

f3 = add_shape(slide, Inches(9), facts_y, Inches(3.6), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B), ACCENT_GREEN)
add_text_box(slide, Inches(9.3), Inches(4.85), Inches(3), Inches(0.4),
             "Prompt = Instruktion", font_size=18, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(9.3), Inches(5.35), Inches(3), Inches(1.2),
             "Du skriver en prompt (instruktion) och modellen genererar text som svar",
             font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5: Vad kan man göra med Generativ AI?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Vad kan man göra med Generativ AI?", font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Tre huvudkategorier av uppgifter", font_size=20, color=MEDIUM_GRAY)

# Three columns
col_w = Inches(3.7)
col_h = Inches(4.8)
col_y = Inches(2.0)

# Writing
w_box = add_shape(slide, Inches(0.8), col_y, col_w, col_h, RGBColor(0x1E, 0x29, 0x3B), ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(3.2), Inches(0.5),
             "SKRIVA", font_size=26, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.1), Inches(2.9), Inches(3.2), Inches(3.5),
             "Kort prompt --> Lång text\n\n"
             "- Brainstorming & idéer\n"
             "- Marknadsföringstexter\n"
             "- E-post & rapporter\n"
             "- Översättning\n"
             "- Kodgenerering\n"
             "- Kreativt skrivande",
             font_size=15, color=LIGHT_GRAY)

# Reading
r_box = add_shape(slide, Inches(4.9), col_y, col_w, col_h, RGBColor(0x1E, 0x29, 0x3B), ACCENT_PURPLE)
add_text_box(slide, Inches(5.2), Inches(2.2), Inches(3.2), Inches(0.5),
             "LÄSA", font_size=26, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.2), Inches(2.9), Inches(3.2), Inches(3.5),
             "Lång text --> Kort svar\n\n"
             "- Sammanfattningar\n"
             "- Korrekturläsning\n"
             "- Sentimentanalys\n"
             "- E-postkategorisering\n"
             "- Datautvinning\n"
             "- Klassificering",
             font_size=15, color=LIGHT_GRAY)

# Chatting
c_box = add_shape(slide, Inches(9), col_y, col_w, col_h, RGBColor(0x1E, 0x29, 0x3B), ACCENT_GREEN)
add_text_box(slide, Inches(9.3), Inches(2.2), Inches(3.2), Inches(0.5),
             "CHATTA", font_size=26, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.3), Inches(2.9), Inches(3.2), Inches(3.5),
             "Interaktiv konversation\n\n"
             "- Kundtjänst-chatbotar\n"
             "- IT-support\n"
             "- Karriärrådgivning\n"
             "- Beställningssystem\n"
             "- Utbildningshjälp\n"
             "- Specialiserade assistenter",
             font_size=15, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6: Promptteknik
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Promptteknik - Konsten att prata med AI", font_size=38, color=WHITE, bold=True)

# Tip 1
t1 = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.3), RGBColor(0x1E, 0x29, 0x3B), ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5.3), Inches(0.4),
             "1. Var detaljerad och specifik", font_size=20, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.1), Inches(5.3), Inches(0.5),
             'Dåligt: "Skriv ett mail"', font_size=14, color=ACCENT_PINK)
add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5.3), Inches(1),
             'Bra: "Skriv ett professionellt mail till min lärare\nom att jag behöver mer tid för mitt projekt.\nTonen ska vara artig men tydlig."',
             font_size=14, color=ACCENT_GREEN)

# Tip 2
t2 = add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(2.3), RGBColor(0x1E, 0x29, 0x3B), ACCENT_PURPLE)
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.3), Inches(0.4),
             "2. Ge steg-för-steg instruktioner", font_size=20, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5.3), Inches(1.5),
             "Steg 1: Generera 5 idéer\nSteg 2: Välj den bästa\nSteg 3: Utveckla den med detaljer\nSteg 4: Skriv en kort sammanfattning",
             font_size=14, color=LIGHT_GRAY)

# Tip 3
t3 = add_shape(slide, Inches(0.8), Inches(4.1), Inches(5.8), Inches(2.8), RGBColor(0x1E, 0x29, 0x3B), ACCENT_CYAN)
add_text_box(slide, Inches(1.1), Inches(4.2), Inches(5.3), Inches(0.4),
             "3. Iterera och förbättra", font_size=20, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(1.1), Inches(4.7), Inches(5.3), Inches(2),
             "Prompting är en loop:\n\n"
             "Idé --> Prompt --> Resultat --> Förbättra --> Upprepa\n\n"
             "Det första försöket behöver inte vara perfekt!\nBörja enkelt och förfina stegvis.",
             font_size=14, color=LIGHT_GRAY)

# Warning box
w_box = add_shape(slide, Inches(7), Inches(4.1), Inches(5.8), Inches(2.8), RGBColor(0x2D, 0x1B, 0x1B), ACCENT_ORANGE)
add_text_box(slide, Inches(7.3), Inches(4.2), Inches(5.3), Inches(0.4),
             "Viktigt att tänka på!", font_size=20, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(7.3), Inches(4.7), Inches(5.3), Inches(2),
             "- Dela aldrig känslig/personlig information\n  med AI-verktyg\n\n"
             "- Dubbelkolla alltid viktiga fakta\n  (AI kan hallucinera)\n\n"
             "- Speciellt viktigt för medicinska,\n  juridiska och ekonomiska beslut",
             font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 7: Begränsningar
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Begränsningar - Vad AI INTE kan", font_size=40, color=WHITE, bold=True)

limits = [
    ("Kunskapsgräns", "Modellen vet bara det den tränats på. Den har ingen tillgång till aktuella händelser efter sitt träningsdatum.", ACCENT_BLUE),
    ("Hallucinationer", "AI:n kan hitta på saker som låter helt trovärdiga men är helt fel. Kontrollera alltid viktiga fakta!", ACCENT_PINK),
    ("Begränsad input", "Det finns gränser för hur mycket text modellen kan bearbeta åt gången (kontextfönster).", ACCENT_PURPLE),
    ("Svag med tabeller", "Strukturerad data som Excel-tabeller hanteras bättre med traditionella verktyg.", ACCENT_ORANGE),
    ("Bias & fördomar", "Modellen kan spegla fördomar från sin träningsdata och ge partiska eller stereotypa svar.", ACCENT_CYAN),
]

for i, (title, desc, color) in enumerate(limits):
    y = Inches(1.5 + i * 1.15)
    box = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.0), RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(3), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(4), y + Inches(0.05), Inches(8), Inches(0.85),
                 desc, font_size=15, color=LIGHT_GRAY)


# ============================================================
# SLIDE 8: Bildgenerering
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Bildgenerering med AI", font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "Diffusionsmodeller - från brus till bilder", font_size=20, color=ACCENT_CYAN)

# Process illustration
steps = [
    ("1", "Brus", "Börja med\nslumpmässiga pixlar", ACCENT_BLUE),
    ("2", "Rensa", "Modellen tar bort\nlite brus i taget", ACCENT_PURPLE),
    ("3", "Form", "Konturer och\nformer framträder", ACCENT_CYAN),
    ("4", "Detalj", "Detaljer och\nfärger läggs till", ACCENT_GREEN),
    ("5", "Resultat", "En färdig realistisk\nbild skapas", ACCENT_ORANGE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.8 + i * 2.5)
    box = add_shape(slide, x, Inches(2.0), Inches(2.2), Inches(2.5), RGBColor(0x1E, 0x29, 0x3B), color)
    # Step number
    circle = add_circle(slide, x + Inches(0.7), Inches(2.2), Inches(0.7), color)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.1), Inches(3.1), Inches(2), Inches(0.4),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.5), Inches(2), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.2), Inches(3.0), Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MEDIUM_GRAY
        arrow.line.fill.background()

# Prompt control box
prompt_box = add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2), RGBColor(0x1E, 0x29, 0x3B), ACCENT_PURPLE)
add_text_box(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.4),
             "Textprompts styr resultatet", font_size=20, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(1.2),
             'Du kan styra bildgenereringen med textbeskrivningar:\n\n'
             'Prompt: "En futuristisk stad vid solnedgång, cyberpunk-stil" --> AI genererar en bild som matchar beskrivningen\n'
             'Prompt: "En tecknad katt som spelar gitarr" --> AI genererar en helt annan bild baserat på den nya beskrivningen\n\n'
             'Verktyg: DALL-E, Midjourney, Stable Diffusion',
             font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 9: Verkliga användningsområden
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Verkliga användningsområden", font_size=40, color=WHITE, bold=True)

areas = [
    ("Utbildning", "Personlig handledning, studiehjälp, skapa övningar, förklara komplexa ämnen", ACCENT_BLUE),
    ("Hälsa", "Medicinsk forskning, analysera patientdata, hjälpa läkare med diagnostik", ACCENT_GREEN),
    ("Kreativitet", "Musik, konst, design, filmproduktion, speldesign, arkitektur", ACCENT_PURPLE),
    ("Företag", "Kundtjänst, marknadsföring, dataanalys, automatisering, rapportskrivning", ACCENT_ORANGE),
    ("Utveckling", "Kodgenerering, buggfixning, dokumentation, testning, code review", ACCENT_CYAN),
    ("Vetenskap", "Proteinveckning, läkemedelsutveckling, klimatmodellering, materialforskning", ACCENT_PINK),
]

for i, (title, desc, color) in enumerate(areas):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.5 + row * 2.8)

    box = add_shape(slide, x, y, Inches(3.7), Inches(2.4), RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5),
                 title, font_size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.4),
                 desc, font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 10: Kursen - Vad ska vi göra?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Kursen - Vad ska vi göra?", font_size=40, color=WHITE, bold=True)

# Three phases
phases = [
    ("Del 1", "Teori & Grunder", "Förstå vad generativ AI är, hur det fungerar, och vad det kan användas till. Lär dig promptteknik och kritiskt tänkande kring AI.", ACCENT_BLUE, "Vecka 1-2"),
    ("Del 2", "Praktik & Kursmaterial", "Arbeta igenom Microsofts kursmaterial 'Generative AI with JavaScript'. Bygga AI-applikationer med JavaScript och API:er.", ACCENT_PURPLE, "Vecka 3-6"),
    ("Del 3", "Eget Projekt", "Bygg din egen AI-applikation! Använd det du lärt dig för att skapa något unikt. Presentera för klassen.", ACCENT_GREEN, "Vecka 7-10"),
]

for i, (part, title, desc, color, time) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    box = add_shape(slide, x, Inches(1.5), Inches(3.7), Inches(5.2), RGBColor(0x1E, 0x29, 0x3B), color)

    # Part label
    label_box = add_shape(slide, x + Inches(0.2), Inches(1.7), Inches(1.2), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.25), Inches(1.72), Inches(1.1), Inches(0.45),
                 part, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.5),
                 title, font_size=22, color=color, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.3), Inches(3.1), Inches(3.2), Inches(2.5),
                 desc, font_size=15, color=LIGHT_GRAY)

    # Timeline
    add_text_box(slide, x + Inches(0.3), Inches(5.8), Inches(3.2), Inches(0.4),
                 time, font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: Projektidéer
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Projektidéer - Vad kan du bygga?", font_size=40, color=WHITE, bold=True)

ideas = [
    ("Studiehjälp", "En AI som hjälper dig plugga genom att generera frågor och förklaringar", ACCENT_BLUE),
    ("Chatbot", "En specialiserad chatbot för ett ämne du brinner för", ACCENT_PURPLE),
    ("Nyhetssammanfattare", "En app som sammanfattar långa nyhetsartiklar", ACCENT_CYAN),
    ("Kreativ assistent", "Ett verktyg för att skapa berättelser, dikter eller låttexter", ACCENT_GREEN),
    ("Kodhjälpare", "En AI-partner som hjälper dig lära dig programmera", ACCENT_ORANGE),
    ("Bildgenerator", "En app där du beskriver bilder med text och AI skapar dem", ACCENT_PINK),
    ("Receptgenerator", "Beskriv vad du har i kylen och få ett recept!", ACCENT_BLUE),
    ("Språkträning", "En konversationspartner för att öva ett nytt språk", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(ideas):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(1.5 + row * 2.8)

    box = add_shape(slide, x, y, Inches(2.8), Inches(2.4), RGBColor(0x1E, 0x29, 0x3B), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.5), Inches(1.4),
                 desc, font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 12: Avslutning
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative shapes
add_circle(slide, Inches(9), Inches(-1), Inches(4), ACCENT_PURPLE)
add_circle(slide, Inches(-1), Inches(4), Inches(5), ACCENT_BLUE)

# Overlay
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
spPr2 = overlay._element.find(f'.//{{{a_ns}}}solidFill/{{{a_ns}}}srgbClr')
if spPr2 is not None:
    alpha2 = etree.SubElement(spPr2, f'{{{a_ns}}}alpha')
    alpha2.set('val', '75000')
overlay.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Redo att börja?", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
             "Generativ AI förändrar världen.\nVi ska lära oss att använda det - och bygga med det.",
             font_size=24, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Divider
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.3), Inches(3.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.6),
             "github.com/microsoft/generative-ai-with-javascript", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
             "Låt oss skapa framtiden tillsammans!", font_size=28, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)


# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Generativ-AI-Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
